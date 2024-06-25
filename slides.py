from pptx import Presentation
from pptx.util import Inches, Pt

def create_slide(prs, layout, title, content):
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = content

def main():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # Title Slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Soaring to New Heights"
    subtitle.text = "Exploring the Luxurious Performance of the Cirrus SR22T"

    # Introduction Slide
    intro_content = (
        "• Experience first-class flying at 17,000 feet\n"
        "• Redefining private aviation\n"
        "• Combining luxury and performance"
    )
    create_slide(prs, content_slide_layout, "Introduction", intro_content)

    # Performance Specifications
    specs_content = (
        "• Max Cruise Speed: 213 ktas\n"
        "• Range: 1,021 nm\n"
        "• Useful Load: 1,328 lbs\n"
        "• Service Ceiling: 25,000 ft\n"
        "• Takeoff Distance: 1,517 ft"
    )
    create_slide(prs, content_slide_layout, "Performance Specifications", specs_content)

    # Luxury Features
    luxury_content = (
        "• Premium leather seats\n"
        "• Cirrus Perspective+ by Garmin® flight deck\n"
        "• Spacious cabin with panoramic windows\n"
        "• Climate control system\n"
        "• Noise-reducing headsets"
    )
    create_slide(prs, content_slide_layout, "Luxury Features", luxury_content)

    # Safety Innovations
    safety_content = (
        "• Cirrus Airframe Parachute System® (CAPS®)\n"
        "• Enhanced Vision System (EVS)\n"
        "• Electronic Stability & Protection (ESP)\n"
        "• Hypoxia Recognition System with Autopilot Descent\n"
        "• Airbag seatbelts"
    )
    create_slide(prs, content_slide_layout, "Safety Innovations", safety_content)

    # Conclusion
    conclusion_content = (
        "• The Cirrus SR22T offers:\n"
        "  - Unparalleled luxury in private aviation\n"
        "  - Impressive performance capabilities\n"
        "  - Cutting-edge safety features\n"
        "• Elevate your flying experience with the Cirrus SR22T"
    )
    create_slide(prs, content_slide_layout, "Conclusion", conclusion_content)

    prs.save('Cirrus_SR22T_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    main()